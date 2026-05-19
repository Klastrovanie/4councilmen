"""
document_parser.py
==================
Temporary upload storage and document extraction helpers for 4Councilmen.

Supported inputs:
- PDF: text extraction via pdfplumber / pypdf, OCR fallback for scanned pages when available
- CSV / XLSX / XLS: pandas table summaries
- DOCX / DOC: python-docx, DOC via LibreOffice conversion when available
- TXT: plain text
- PPTX / PPT: slide text + speaker notes + embedded-image OCR, PPT via LibreOffice conversion when available
- JPG / JPEG / PNG: OCR when available

Files are stored under /app/tmp_uploads/<session_id>/ and are intended to be
removed after one 4CM run. All extraction is best-effort and bounded to control
API cost and prompt length.
"""

from __future__ import annotations

import csv
import hashlib
import os
import re
import shutil
import subprocess
import time
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional, Tuple

UPLOAD_ROOT = Path(os.environ.get("FOURCM_UPLOAD_ROOT", "/app/tmp_uploads"))
MAX_FILES = int(os.environ.get("FOURCM_MAX_FILES", "10"))
MAX_FILE_BYTES = int(os.environ.get("FOURCM_MAX_FILE_BYTES", str(25 * 1024 * 1024)))
MAX_TOTAL_BYTES = int(os.environ.get("FOURCM_MAX_TOTAL_BYTES", str(100 * 1024 * 1024)))
MAX_CONTEXT_CHARS = int(os.environ.get("FOURCM_MAX_CONTEXT_CHARS", "60000"))
MAX_CHARS_PER_FILE = int(os.environ.get("FOURCM_MAX_CHARS_PER_FILE", "12000"))
MAX_PDF_PAGES = int(os.environ.get("FOURCM_MAX_PDF_PAGES", "30"))
MAX_OCR_PAGES = int(os.environ.get("FOURCM_MAX_OCR_PAGES", "8"))
MAX_OCR_IMAGES = int(os.environ.get("FOURCM_MAX_OCR_IMAGES", "20"))
UPLOAD_TTL_SECONDS = int(os.environ.get("FOURCM_UPLOAD_TTL_SECONDS", "3600"))

ALLOWED_EXTENSIONS = {
    ".pdf", ".csv", ".xlsx", ".xls", ".docx", ".doc", ".txt",
    ".pptx", ".ppt", ".jpg", ".jpeg", ".png",
}

TEXT_EXTENSIONS = {".txt"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png"}


@dataclass
class StoredUpload:
    original_name: str
    stored_name: str
    size: int
    extension: str


def ensure_upload_root() -> None:
    UPLOAD_ROOT.mkdir(parents=True, exist_ok=True)


def new_session_id() -> str:
    return f"{int(time.time())}_{uuid.uuid4().hex[:16]}"


def session_path(session_id: str) -> Path:
    if not re.fullmatch(r"[0-9A-Za-z_.-]+", session_id or ""):
        raise ValueError("Invalid upload session id")
    return UPLOAD_ROOT / session_id


def sanitize_filename(filename: str) -> str:
    name = Path(filename or "uploaded_file").name
    name = re.sub(r"[^0-9A-Za-z가-힣._ -]+", "_", name).strip(" .")
    if not name:
        name = "uploaded_file"
    stem = Path(name).stem[:80] or "uploaded_file"
    suffix = Path(name).suffix.lower()[:16]
    return f"{stem}{suffix}"


def validate_extension(filename: str) -> str:
    ext = Path(filename or "").suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext or '(none)'}")
    return ext


def cleanup_upload_session(session_id: Optional[str]) -> None:
    if not session_id:
        return
    try:
        path = session_path(session_id)
        if path.exists() and path.is_dir():
            shutil.rmtree(path, ignore_errors=True)
    except Exception:
        # Cleanup must never break the response stream.
        pass


def cleanup_old_upload_sessions(max_age_seconds: int = UPLOAD_TTL_SECONDS) -> None:
    ensure_upload_root()
    now = time.time()
    for p in UPLOAD_ROOT.iterdir():
        try:
            if p.is_dir() and (now - p.stat().st_mtime) > max_age_seconds:
                shutil.rmtree(p, ignore_errors=True)
        except Exception:
            pass


def list_session_files(session_id: str) -> List[Path]:
    path = session_path(session_id)
    if not path.exists() or not path.is_dir():
        return []
    return sorted([p for p in path.iterdir() if p.is_file()])


def _truncate(text: str, limit: int) -> str:
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n\n[TRUNCATED after {limit:,} characters to control AI API cost.]"


def _clean_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _import_or_none(module_name: str):
    try:
        return __import__(module_name)
    except Exception:
        return None


def _ocr_image(image_path: Path) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(image_path)
        # Korean+English when kor data is installed; fallback to English.
        try:
            return pytesseract.image_to_string(img, lang="kor+eng")
        except Exception:
            return pytesseract.image_to_string(img, lang="eng")
    except Exception as e:
        return f"[OCR unavailable or failed for image: {e}]"


def _extract_txt(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp949", "latin-1"):
        try:
            return path.read_text(encoding=enc, errors="replace")
        except Exception:
            continue
    return path.read_bytes().decode("utf-8", errors="replace")


def _extract_pdf(path: Path) -> str:
    parts: List[str] = []
    extracted_any = False

    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            for idx, page in enumerate(pdf.pages[:MAX_PDF_PAGES], start=1):
                try:
                    text = page.extract_text() or ""
                    if text.strip():
                        extracted_any = True
                        parts.append(f"[Page {idx}]\n{text}")
                    # Table preview, bounded.
                    tables = page.extract_tables() or []
                    for t_i, table in enumerate(tables[:3], start=1):
                        rows = table[:12]
                        table_text = "\n".join(" | ".join(str(c or "") for c in row) for row in rows)
                        if table_text.strip():
                            parts.append(f"[Page {idx} table {t_i} preview]\n{table_text}")
                except Exception:
                    continue
    except Exception:
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            for idx, page in enumerate(reader.pages[:MAX_PDF_PAGES], start=1):
                text = page.extract_text() or ""
                if text.strip():
                    extracted_any = True
                    parts.append(f"[Page {idx}]\n{text}")
        except Exception as e:
            parts.append(f"[PDF text extraction failed: {e}]")

    # OCR fallback for scanned/mostly image PDFs.
    if not extracted_any or len("\n".join(parts)) < 800:
        try:
            from pdf2image import convert_from_path
            pages = convert_from_path(str(path), first_page=1, last_page=MAX_OCR_PAGES, dpi=180)
            ocr_parts = []
            tmpdir = path.parent / f".__ocr_{path.stem}"
            tmpdir.mkdir(exist_ok=True)
            for i, img in enumerate(pages, start=1):
                img_path = tmpdir / f"page_{i}.png"
                img.save(img_path)
                ocr_text = _ocr_image(img_path)
                if ocr_text.strip():
                    ocr_parts.append(f"[OCR page {i}]\n{ocr_text}")
            shutil.rmtree(tmpdir, ignore_errors=True)
            if ocr_parts:
                parts.append("\n".join(ocr_parts))
        except Exception as e:
            parts.append(f"[PDF OCR skipped/unavailable: {e}]")

    return "\n\n".join(parts)


def _safe_read_dataframe(path: Path):
    import pandas as pd
    ext = path.suffix.lower()
    if ext == ".csv":
        # Try common encodings.
        for enc in ("utf-8", "utf-8-sig", "cp949", "latin-1"):
            try:
                return pd.read_csv(path, encoding=enc)
            except Exception:
                continue
        return pd.read_csv(path, encoding="utf-8", errors="replace")
    # xlsx/xls
    return pd.read_excel(path, sheet_name=None)


def _summarize_dataframe(df, sheet_name: str = "") -> str:
    import pandas as pd
    label = f"Sheet: {sheet_name}" if sheet_name else "Table"
    parts = [f"[{label}] rows={len(df)}, columns={len(df.columns)}"]
    parts.append("Columns: " + ", ".join(map(str, df.columns[:80])))
    try:
        preview = df.head(12).to_csv(index=False)
        parts.append("Preview (first rows):\n" + preview)
    except Exception:
        pass

    numeric = df.select_dtypes(include="number")
    if not numeric.empty:
        try:
            desc = numeric.describe().transpose().round(4).to_csv()
            parts.append("Numeric summary:\n" + desc)
        except Exception:
            pass

    # Simple finance-oriented risk flags based on column names.
    flags = []
    joined_cols = " ".join(str(c).lower() for c in df.columns)
    if any(k in joined_cols for k in ["revenue", "sales", "매출"]):
        flags.append("Revenue/sales fields detected; compare growth against cash-flow and receivables.")
    if any(k in joined_cols for k in ["receivable", "accounts receivable", "매출채권"]):
        flags.append("Receivables fields detected; check whether receivables grow faster than revenue.")
    if any(k in joined_cols for k in ["inventory", "재고"]):
        flags.append("Inventory fields detected; check inventory build-up and write-down risk.")
    if any(k in joined_cols for k in ["cash flow", "operating cash", "영업현금"]):
        flags.append("Cash-flow fields detected; compare net income and operating cash flow quality.")
    if flags:
        parts.append("Financial analysis hints:\n- " + "\n- ".join(flags))
    return "\n".join(parts)


def _extract_spreadsheet(path: Path) -> str:
    try:
        data = _safe_read_dataframe(path)
        if isinstance(data, dict):
            parts = []
            for name, df in list(data.items())[:10]:
                parts.append(_summarize_dataframe(df, str(name)))
            return "\n\n".join(parts)
        return _summarize_dataframe(data)
    except Exception as e:
        return f"[Spreadsheet extraction failed: {e}]"


def _extract_docx(path: Path) -> str:
    try:
        import docx
        doc = docx.Document(str(path))
        parts = []
        paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        if paras:
            parts.append("\n".join(paras))
        for ti, table in enumerate(doc.tables[:20], start=1):
            rows = []
            for row in table.rows[:20]:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            if rows:
                parts.append(f"[Table {ti}]\n" + "\n".join(rows))
        return "\n\n".join(parts)
    except Exception as e:
        return f"[DOCX extraction failed: {e}]"


def _convert_with_libreoffice(path: Path, target_ext: str) -> Optional[Path]:
    outdir = path.parent / f".__converted_{path.stem}"
    outdir.mkdir(exist_ok=True)
    try:
        cmd = [
            "libreoffice", "--headless", "--convert-to", target_ext.lstrip("."),
            "--outdir", str(outdir), str(path)
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=90)
        converted = outdir / f"{path.stem}{target_ext}"
        if converted.exists():
            return converted
        matches = list(outdir.glob(f"*{target_ext}"))
        return matches[0] if matches else None
    except Exception:
        shutil.rmtree(outdir, ignore_errors=True)
        return None


def _extract_legacy_office(path: Path) -> str:
    ext = path.suffix.lower()
    target = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx"}.get(ext)
    if not target:
        return "[Unsupported legacy Office file]"
    converted = _convert_with_libreoffice(path, target)
    if not converted:
        return "[Legacy Office conversion unavailable. Install LibreOffice in the Docker image.]"
    try:
        return extract_file_text(converted)
    finally:
        shutil.rmtree(converted.parent, ignore_errors=True)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        image_count = 0
        tmpdir = path.parent / f".__pptx_images_{path.stem}"
        tmpdir.mkdir(exist_ok=True)
        for si, slide in enumerate(prs.slides[:80], start=1):
            slide_parts = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                except Exception:
                    pass
                try:
                    if image_count < MAX_OCR_IMAGES and getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
                        image = shape.image
                        image_count += 1
                        img_path = tmpdir / f"slide_{si}_image_{image_count}.{image.ext}"
                        img_path.write_bytes(image.blob)
                        ocr = _ocr_image(img_path)
                        if ocr.strip():
                            slide_parts.append(f"[OCR image {image_count}]\n{ocr}")
                except Exception:
                    pass
            if slide_parts:
                parts.append(f"[Slide {si}]\n" + "\n".join(slide_parts))
        shutil.rmtree(tmpdir, ignore_errors=True)
        return "\n\n".join(parts)
    except Exception as e:
        return f"[PPTX extraction failed: {e}]"


def _extract_image(path: Path) -> str:
    return _ocr_image(path)


def extract_file_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in TEXT_EXTENSIONS:
        return _clean_text(_extract_txt(path))
    if ext == ".pdf":
        return _clean_text(_extract_pdf(path))
    if ext in {".csv", ".xlsx"}:
        return _clean_text(_extract_spreadsheet(path))
    if ext == ".xls":
        return _clean_text(_extract_legacy_office(path))
    if ext == ".docx":
        return _clean_text(_extract_docx(path))
    if ext == ".doc":
        return _clean_text(_extract_legacy_office(path))
    if ext == ".pptx":
        return _clean_text(_extract_pptx(path))
    if ext == ".ppt":
        return _clean_text(_extract_legacy_office(path))
    if ext in IMAGE_EXTENSIONS:
        return _clean_text(_extract_image(path))
    return f"[Unsupported file type: {ext}]"


def build_document_context(session_id: str) -> str:
    files = list_session_files(session_id)
    if not files:
        return ""

    chunks: List[str] = []
    total = 0
    for idx, path in enumerate(files[:MAX_FILES], start=1):
        try:
            digest = hashlib.sha256(path.read_bytes()).hexdigest()[:16]
        except Exception:
            digest = "unknown"
        extracted = extract_file_text(path)
        extracted = _truncate(extracted, MAX_CHARS_PER_FILE)
        item = (
            f"[Uploaded File {idx}: {path.name}]\n"
            f"- extension: {path.suffix.lower()}\n"
            f"- size_bytes: {path.stat().st_size if path.exists() else 0}\n"
            f"- sha256_prefix: {digest}\n"
            f"- extracted_text_or_table_summary:\n{extracted or '[No readable text extracted]'}"
        )
        chunks.append(item)
        total += len(item)
        if total >= MAX_CONTEXT_CHARS:
            chunks.append(f"[DOCUMENT CONTEXT TRUNCATED after {MAX_CONTEXT_CHARS:,} characters to control AI API cost.]")
            break

    return "\n\n---\n\n".join(chunks)
